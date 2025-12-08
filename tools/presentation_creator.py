import logging
import os
import json
from typing import List, Dict, Any, Union

from tools.base_tool import BaseTool
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

logger = logging.getLogger(__name__)

class PresentationCreatorTool(BaseTool):
    """
    A tool for creating presentations in various formats (PPTX, PDF, HTML, Markdown)
    from structured slide data.
    """

    def __init__(self, tool_name: str = "PresentationCreator", data_dir: str = ".", **kwargs):
        super().__init__(tool_name=tool_name, **kwargs)
        self.output_dir = os.path.join(data_dir, "presentations")
        os.makedirs(self.output_dir, exist_ok=True)

    @property
    def description(self) -> str:
        return "Creates presentations in PPTX, PDF, HTML, or Markdown formats from structured slide data."

    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "operation": {"type": "string", "enum": ["create_presentation"]},
                "file_name": {"type": "string", "description": "The base name for the output file (e.g., 'MyPresentation')."},
                "format": {"type": "string", "enum": ["pptx", "pdf", "html", "md"], "default": "pptx"},
                "slides_data": {
                    "type": "array",
                    "items": {"type": "object", "properties": {"title": {"type": "string"}, "content": {"type": "string"}}},
                    "description": "A list of dictionaries, each representing a slide with 'title' and 'content'."
                }
            },
            "required": ["operation", "file_name", "slides_data"]
        }

    def _create_pptx(self, file_path: str, slides_data: List[Dict[str, Any]]) -> str:
        """Creates a PowerPoint presentation (.pptx) from structured data."""
        prs = Presentation()
        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data.get("title", "No Title")
            
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.text = slide_data.get("content", "")

        prs.save(file_path)
        return f"Successfully generated PowerPoint presentation at '{file_path}'."

    def _create_pdf(self, file_path: str, slides_data: List[Dict[str, Any]]) -> str:
        """Creates a basic PDF presentation from structured data."""
        c = canvas.Canvas(file_path, pagesize=letter)
        for slide_data in slides_data:
            y_position = 750
            c.setFont("Helvetica-Bold", 24)
            c.drawString(100, y_position, slide_data.get("title", ""))
            y_position -= 30
            c.setFont("Helvetica", 12)
            textobject = c.beginText()
            textobject.setTextOrigin(100, y_position)
            for line in slide_data.get("content", "").split('\n'):
                textobject.textLine(line)
            c.drawText(textobject)
            c.showPage() # New page for each slide
        c.save()
        return f"Successfully generated PDF presentation at '{file_path}'."

    def _create_html(self, file_path: str, slides_data: List[Dict[str, Any]]) -> str:
        """Creates an HTML presentation from structured data."""
        html_content = """
<!DOCTYPE html>
<html>
<head>
    <title>Presentation</title>
    <style>
        body { font-family: sans-serif; margin: 0; padding: 0; background-color: #f0f0f0; }
        .slide { 
            background-color: white; border: 1px solid #ccc; margin: 20px auto; padding: 20px; 
            width: 80%; box-shadow: 2px 2px 5px rgba(0,0,0,0.1);
        }
        h1 { color: #333; border-bottom: 1px solid #eee; padding-bottom: 10px; margin-bottom: 15px; }
        p { line-height: 1.6; color: #555; }
    </style>
</head>
<body>
"""
        for i, slide_data in enumerate(slides_data):
            html_content += f"""
    <div class="slide">
        <h1>{slide_data.get("title", f"Slide {i+1}")}</h1>
        <p>{slide_data.get("content", "")}</p>
    </div>
"""
        html_content += """
</body>
</html>
"""
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        return f"Successfully generated HTML presentation at '{file_path}'."

    def _create_markdown(self, file_path: str, slides_data: List[Dict[str, Any]]) -> str:
        """Creates a Markdown presentation from structured data."""
        md_content = ""
        for i, slide_data in enumerate(slides_data):
            md_content += f"# {slide_data.get("title", f"Slide {i+1}")}\n\n"
            md_content += f"{slide_data.get("content", "")}\n\n---\n\n"
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        return f"Successfully generated Markdown presentation at '{file_path}'."

    def execute(self, operation: str, file_name: str, format: str, slides_data: List[Dict[str, Any]], **kwargs: Any) -> Dict[str, Any]:
        """
        Creates a presentation in the specified format from structured slide data.
        """
        if operation != "create_presentation":
            raise ValueError(f"Invalid operation: {operation}. Only 'create_presentation' is supported.")
        if not slides_data:
            raise ValueError("slides_data cannot be empty.")

        output_file_path = os.path.join(self.output_dir, f"{file_name}.{format}")

        if format == "pptx":
            message = self._create_pptx(output_file_path, slides_data)
        elif format == "pdf":
            message = self._create_pdf(output_file_path, slides_data)
        elif format == "html":
            message = self._create_html(output_file_path, slides_data)
        elif format == "md":
            message = self._create_markdown(output_file_path, slides_data)
        else:
            raise ValueError(f"Unsupported format: {format}. Choose from 'pptx', 'pdf', 'html', 'md'.")

        return {"status": "success", "output_file": output_file_path, "message": message}

if __name__ == '__main__':
    print("Demonstrating PresentationCreatorTool functionality...")
    temp_dir = "temp_presentations"
    if os.path.exists(temp_dir): import shutil; shutil.rmtree(temp_dir)
    os.makedirs(temp_dir)
    
    creator_tool = PresentationCreatorTool(data_dir=temp_dir)
    
    sample_slides = [
        {"title": "Introduction to AI", "content": "Artificial Intelligence is transforming industries."},
        {"title": "Machine Learning Basics", "content": "Supervised, Unsupervised, and Reinforcement Learning."},
        {"title": "Deep Learning", "content": "Neural Networks and their applications."},
        {"title": "Future of AI", "content": "Ethical considerations and societal impact."}
    ]
    
    try:
        # 1. Create a PPTX presentation
        print("\n--- Creating 'AI_Overview.pptx' ---")
        pptx_result = creator_tool.execute(operation="create_presentation", file_name="AI_Overview", format="pptx", slides_data=sample_slides)
        print(json.dumps(pptx_result, indent=2))

        # 2. Create a PDF presentation
        print("\n--- Creating 'AI_Overview.pdf' ---")
        pdf_result = creator_tool.execute(operation="create_presentation", file_name="AI_Overview", format="pdf", slides_data=sample_slides)
        print(json.dumps(pdf_result, indent=2))

        # 3. Create an HTML presentation
        print("\n--- Creating 'AI_Overview.html' ---")
        html_result = creator_tool.execute(operation="create_presentation", file_name="AI_Overview", format="html", slides_data=sample_slides)
        print(json.dumps(html_result, indent=2))

        # 4. Create a Markdown presentation
        print("\n--- Creating 'AI_Overview.md' ---")
        md_result = creator_tool.execute(operation="create_presentation", file_name="AI_Overview", format="md", slides_data=sample_slides)
        print(json.dumps(md_result, indent=2))

    except Exception as e:
        print(f"\nAn error occurred: {e}")
    finally:
        if os.path.exists(temp_dir): import shutil; shutil.rmtree(temp_dir)
        print(f"\nCleaned up temporary directory '{temp_dir}'.")